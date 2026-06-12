#!/usr/bin/env python3
"""把公式 PNG 等比插入 pptx 指定页。绝不拉伸：按 min(w_box/w_img, h_box/h_img) 缩放。

用法:
    python3 insert_formula.py deck.pptx SLIDE_NO formula.png --box X Y MAX_W MAX_H [--align lt|ctr] [-o out.pptx]

SLIDE_NO 从 1 开始；--box 单位英寸：在 (X,Y) 起的 MAX_W x MAX_H 框内等比放置。
--align ctr 时图片在框内居中（默认左上对齐）。
"""
import argparse
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

ap = argparse.ArgumentParser()
ap.add_argument('pptx')
ap.add_argument('slide_no', type=int)
ap.add_argument('image')
ap.add_argument('--box', nargs=4, type=float, required=True, metavar=('X', 'Y', 'MAX_W', 'MAX_H'))
ap.add_argument('--align', choices=['lt', 'ctr'], default='lt')
ap.add_argument('-o', '--out')
a = ap.parse_args()

x, y, mw, mh = a.box
iw, ih = Image.open(a.image).size
s = min(mw / iw, mh / ih)          # 等比缩放系数（像素→英寸）
w, h = iw * s, ih * s
if a.align == 'ctr':
    x, y = x + (mw - w) / 2, y + (mh - h) / 2

prs = Presentation(a.pptx)
slide = list(prs.slides)[a.slide_no - 1]
slide.shapes.add_picture(a.image, Inches(x), Inches(y), Inches(w), Inches(h))
out = a.out or a.pptx
prs.save(out)
print(f'inserted {a.image} -> slide {a.slide_no} at ({x:.2f},{y:.2f}) size {w:.2f}x{h:.2f}in (aspect kept) -> {out}')
